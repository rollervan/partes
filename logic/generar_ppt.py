import io
import json  # Importamos librería para leer el archivo
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def limpiar_texto(texto):
    """Ayuda a normalizar textos para comparaciones"""
    if not texto: return ""
    return texto.strip().lower()

def buscar_diapositiva(prs, titulo_buscado):
    """Busca una slide que contenga el texto del título buscado."""
    titulo_buscado = limpiar_texto(titulo_buscado)
    
    for slide in prs.slides:
        if slide.shapes.title and slide.shapes.title.text:
            titulo_slide = limpiar_texto(slide.shapes.title.text)
            if titulo_buscado in titulo_slide:
                return slide
    return None

def df_to_ppt_table(slide, df, left, top, width, height):
    """Dibuja la tabla en la slide indicada."""
    # Limitamos filas para que no se salga de la diapo
    df_to_print = df.head(12) 
    rows, cols = df_to_print.shape
    
    table_shape = slide.shapes.add_table(rows + 1, cols, left, top, width, height)
    table = table_shape.table

    # Estilos básicos
    for col_idx, col_name in enumerate(df_to_print.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(9)
        p.alignment = PP_ALIGN.CENTER
        # Fondo gris claro para cabecera
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(220, 220, 220)

    for row_idx, row in df_to_print.iterrows():
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            val_str = f"{value:.2f}" if isinstance(value, float) else str(value)
            cell.text = val_str
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(8)
            p.alignment = PP_ALIGN.CENTER

def agregar_placeholder_ia(slide, texto_marcador):
    """Añade un cuadro de texto visual indicando dónde irá la IA."""
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(2)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = f"[{texto_marcador}]"
    p = tf.paragraphs[0]
    p.font.color.rgb = RGBColor(255, 0, 0) # Rojo para que destaque
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER

def reemplazar_marcadores(prs, diccionario_datos):
    """
    Recorre TODA la presentación sustituyendo las claves del diccionario 
    por sus valores, manteniendo el estilo original.
    """
    if not diccionario_datos:
        return

    def sustituir_en_parrafo(paragraph, datos):
        """Función auxiliar para procesar un párrafo."""
        for run in paragraph.runs:
            texto_run = run.text
            for marcador, valor in datos.items():
                if marcador in texto_run:
                    # str(valor) asegura que no falle si el dato es un número
                    run.text = texto_run.replace(marcador, str(valor))

    # Recorremos todas las diapositivas
    for slide in prs.slides:
        for shape in slide.shapes:
            # Caso A: Cuadros de texto
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    sustituir_en_parrafo(paragraph, diccionario_datos)
            
            # Caso B: Tablas
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            sustituir_en_parrafo(paragraph, diccionario_datos)

def generar_ppt(df, lista_figuras, ruta_json=None):
    """
    Carga la plantilla fija y rellena datos en slides específicas.
    Args:
        df: DataFrame para la tabla resumen.
        lista_figuras: Lista de tuplas (titulo, figura_matplotlib).
        ruta_json: (Opcional) Ruta al archivo .json con los datos a sustituir.
    """
    # 1. Cargar plantilla desde la carpeta assets
    ruta_plantilla = "assets/Plantilla_ReunionesCoordinacionFinCuatrimestre.pptx"
    try:
        prs = Presentation(ruta_plantilla)
    except Exception as e:
        # Si falla, creamos una en blanco para no romper la app, pero avisamos
        print(f"Error cargando plantilla: {e}")
        prs = Presentation()

    # --- NUEVO: CARGAR JSON Y SUSTITUIR ---
    if ruta_json:
        try:
            with open(ruta_json, 'r', encoding='utf-8') as f:
                datos_diccionario = json.load(f)
            # Llamamos a la función de reemplazo con los datos cargados
            reemplazar_marcadores(prs, datos_diccionario)
        except Exception as e:
            print(f"Error cargando o procesando el archivo JSON: {e}")

    # --- 1. INSERTAR TABLA EN "RESUMEN DEL CUATRIMESTRE" ---
    slide_resumen = buscar_diapositiva(prs, "Resumen del cuatrimestre")
    if slide_resumen:
        # Insertamos la tabla debajo del título
        df_to_ppt_table(slide_resumen, df, Inches(0.5), Inches(2), Inches(9), Inches(4))
    
    # --- 2. PLACEHOLDERS PARA FUNCIONES FUTURAS (IA) ---
    
    # Slide: Informe sobre los partes de asignatura (Docentes)
    slide_partes_docentes = buscar_diapositiva(prs, "Informe sobre los partes de asignatura de docentes")
    if slide_partes_docentes:
        agregar_placeholder_ia(slide_partes_docentes, "AQUÍ SE GENERARÁ EL RESUMEN IA DE DOCENTES")

    # Slide: Informe sobre los partes de asignatura de delegados
    slide_partes_delegados = buscar_diapositiva(prs, "Informe sobre los partes de asignatura de delegados")
    if slide_partes_delegados:
        agregar_placeholder_ia(slide_partes_delegados, "AQUÍ SE GENERARÁ EL RESUMEN IA DE DELEGADOS")

    # Slide: Coordinación entre asignaturas
    slide_coordinacion = buscar_diapositiva(prs, "Coordinación entre asignaturas")
    if slide_coordinacion:
        agregar_placeholder_ia(slide_coordinacion, "AQUÍ IRÁ LA INFO DE COORDINACIÓN")

    # --- 3. AÑADIR GRÁFICAS ---
    # Las añadimos al final porque python-pptx no permite insertar en medio fácilmente.
    layout_grafica = prs.slide_layouts[1] # Título y objeto
    
    if len(prs.slide_layouts) > 5:
        layout_grafica = prs.slide_layouts[5] # Intentamos usar Solo Título

    for titulo_grafica, figura in lista_figuras:
        slide = prs.slides.add_slide(layout_grafica)
        
        # Título
        if slide.shapes.title:
            slide.shapes.title.text = titulo_grafica
        
        # Imagen
        image_stream = io.BytesIO()
        figura.savefig(image_stream, format='png', bbox_inches='tight', dpi=150)
        image_stream.seek(0)
        
        slide.shapes.add_picture(image_stream, Inches(1), Inches(1.5), width=Inches(8))

    # --- GUARDAR ---
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    
    return output
